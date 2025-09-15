#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Enhanced PDF to PowerPoint Converter
Converts PDFs to editable PowerPoint presentations with advanced text, image, and layout detection
"""

import argparse
import json
import sys
import os
import re
from pathlib import Path
from typing import List, Dict, Tuple, Union, Optional
import traceback
import colorsys
from dataclasses import dataclass

# Ensure UTF-8 encoding for output
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')
if sys.stderr.encoding != 'utf-8':
    sys.stderr.reconfigure(encoding='utf-8')

# PDF and image processing
from pdf2image import convert_from_path
import pytesseract
from PIL import Image, ImageFont, ImageDraw
import cv2
import numpy as np
import fitz  # PyMuPDF for advanced PDF analysis
import io
import base64
from PIL import ImageDraw, ImageFont, ImageFilter

# Configure Tesseract path for Windows
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Try to import EasyOCR, but don't fail if it's not available
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError as e:
    print(f"WARNING: EasyOCR not available: {e}")
    print("INFO: Will use Tesseract OCR only")
    easyocr = None
    EASYOCR_AVAILABLE = False

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.picture import Picture


@dataclass
class TextElement:
    """Represents a text element with all its properties"""
    text: str
    bbox: Tuple[int, int, int, int]  # x1, y1, x2, y2
    font_name: str
    font_size: float
    color: Tuple[int, int, int]  # RGB
    bold: bool
    italic: bool
    confidence: float
    alignment: str  # 'left', 'center', 'right'
    
@dataclass
class ImageElement:
    """Represents an image element"""
    bbox: Tuple[int, int, int, int]
    image_data: np.ndarray
    image_type: str  # 'photo', 'diagram', 'chart', 'logo'
    image_format: str  # 'png', 'jpg', 'bmp'
    original_size: Tuple[int, int]
    transparency: bool
    
@dataclass
class ShapeElement:
    """Represents a shape or drawing element"""
    bbox: Tuple[int, int, int, int]
    shape_type: str  # 'rectangle', 'circle', 'line', 'arrow'
    stroke_color: Tuple[int, int, int]
    fill_color: Optional[Tuple[int, int, int]]
    stroke_width: float

@dataclass
class TableElement:
    """Represents a table structure"""
    bbox: Tuple[int, int, int, int]
    rows: List[List[str]]
    cell_positions: List[List[Tuple[int, int, int, int]]]


class EnhancedPDFToPPTConverter:
    def __init__(self, ocr_engine='tesseract', dpi=200):
        """
        Initialize the enhanced PDF to PowerPoint converter
        
        Args:
            ocr_engine: 'tesseract' or 'easyocr'
            dpi: DPI for PDF to image conversion
        """
        # Check if EasyOCR is requested but not available
        if ocr_engine == 'easyocr' and not EASYOCR_AVAILABLE:
            print("WARNING: EasyOCR requested but not available, falling back to Tesseract")
            ocr_engine = 'tesseract'
        
        self.ocr_engine = ocr_engine
        self.dpi = dpi
        self.reader = None
        
        if ocr_engine == 'easyocr' and EASYOCR_AVAILABLE:
            print("INFO: Initializing EasyOCR...")
            self.reader = easyocr.Reader(['en'])
        else:
            print("INFO: Using Tesseract OCR")
            
        print(f"SUCCESS: Enhanced PDF to PPT Converter initialized with {ocr_engine} at {dpi} DPI")

    def detect_pdf_type(self, pdf_path: str) -> str:
        """
        Detect if PDF is text-based or image-based (scanned)
        
        Returns:
            'text-based' or 'image-based'
        """
        try:
            doc = fitz.open(pdf_path)
            total_text_length = 0
            total_pages = len(doc)
            
            for page_num in range(min(3, total_pages)):  # Check first 3 pages
                page = doc[page_num]
                text = page.get_text()
                total_text_length += len(text.strip())
            
            doc.close()
            
            # If we found substantial text, it's likely text-based
            avg_text_per_page = total_text_length / min(3, total_pages)
            
            if avg_text_per_page > 100:  # Threshold for text-based
                return 'text-based'
            else:
                return 'image-based'
                
        except Exception as e:
            print(f"Warning: Could not analyze PDF type: {e}")
            return 'image-based'  # Default to image-based for safety

    def extract_text_elements_with_formatting(self, pdf_path: str) -> List[Dict]:
        """
        Extract text elements with detailed formatting information
        """
        try:
            print(f"INFO: Extracting text elements with formatting from: {pdf_path}")
            
            # First try to extract text with formatting using PyMuPDF
            doc = fitz.open(pdf_path)
            pages_data = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_dict = page.get_text("dict")
                page_data = self._analyze_page_structure(page_dict, page_num + 1)
                pages_data.append(page_data)
            
            doc.close()
            
            # If text extraction was insufficient, use OCR
            for page_data in pages_data:
                if len(page_data['text_elements']) < 3:  # Too little text found
                    print(f"INFO: Supplementing with OCR for page {page_data['page_number']}")
                    ocr_elements = self._extract_with_ocr(pdf_path, page_data['page_number'])
                    page_data['text_elements'].extend(ocr_elements)
            
            return pages_data
            
        except Exception as e:
            print(f"ERROR: Error extracting text elements: {e}")
            traceback.print_exc()
            # Fallback to OCR-only
            return self._extract_with_ocr_only(pdf_path)

    def extract_images_from_pdf(self, pdf_path: str) -> Dict[int, List[ImageElement]]:
        """
        Extract all images from PDF with high quality
        """
        try:
            doc = fitz.open(pdf_path)
            page_images = {}
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_images[page_num + 1] = []
                
                # Get image list
                image_list = page.get_images(full=True)
                
                for img_index, img in enumerate(image_list):
                    try:
                        # Extract image
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        # Convert to PIL Image
                        image = Image.open(io.BytesIO(image_bytes))
                        image_array = np.array(image)
                        
                        # Get image position on page
                        image_rects = page.get_image_rects(img)
                        
                        for rect in image_rects:
                            bbox = (int(rect.x0), int(rect.y0), int(rect.x1), int(rect.y1))
                            
                            image_element = ImageElement(
                                bbox=bbox,
                                image_data=image_array,
                                image_type=self._classify_image(image_array),
                                image_format=image_ext,
                                original_size=image.size,
                                transparency=image.mode in ['RGBA', 'LA']
                            )
                            
                            page_images[page_num + 1].append(image_element)
                            
                    except Exception as e:
                        print(f"WARNING: Could not extract image {img_index} from page {page_num + 1}: {e}")
                        continue
            
            doc.close()
            return page_images
            
        except Exception as e:
            print(f"ERROR: Could not extract images from PDF: {e}")
            return {}

    def _classify_image(self, image_array: np.ndarray) -> str:
        """
        Classify image type based on content analysis
        """
        if len(image_array.shape) == 2 or (len(image_array.shape) == 3 and image_array.shape[2] == 1):
            return 'diagram'  # Grayscale, likely diagram
        
        # Analyze color distribution
        if len(image_array.shape) == 3:
            # Check if it's mostly text (high contrast, limited colors)
            unique_colors = len(np.unique(image_array.reshape(-1, image_array.shape[2]), axis=0))
            if unique_colors < 10:
                return 'diagram'
            elif unique_colors < 50:
                return 'chart'
            else:
                return 'photo'
        
        return 'unknown'

    def _analyze_page_structure(self, page_dict: dict, page_number: int) -> dict:
        """
        Analyze page structure from PyMuPDF dictionary
        """
        text_elements = []
        image_elements = []
        shapes = []
        tables = []
        
        page_width = page_dict.get('width', 612)
        page_height = page_dict.get('height', 792)
        
        # Extract text blocks
        for block in page_dict.get('blocks', []):
            if 'lines' in block:  # Text block
                for line in block['lines']:
                    for span in line['spans']:
                        text = span['text'].strip()
                        if text:
                            bbox = span['bbox']
                            font_info = span.get('font', 'Arial')
                            font_size = span.get('size', 12)
                            font_flags = span.get('flags', 0)
                            color = span.get('color', 0)
                            
                            # Convert color from integer to RGB
                            rgb_color = self._int_to_rgb(color)
                            
                            # Parse font flags (more comprehensive)
                            bold = bool(font_flags & 2**4)
                            italic = bool(font_flags & 2**1)
                            
                            # Extract actual font name and clean it
                            clean_font_name = self._clean_font_name(font_info)
                            
                            # Determine alignment based on position and text analysis
                            alignment = self._determine_alignment_advanced(bbox, page_width, text)
                            
                            text_element = TextElement(
                                text=text,
                                bbox=tuple(bbox),
                                font_name=clean_font_name,
                                font_size=font_size,
                                color=rgb_color,
                                bold=bold,
                                italic=italic,
                                confidence=1.0,  # High confidence for native text
                                alignment=alignment
                            )
                            text_elements.append(text_element)
            
            elif 'ext' in block:  # Image block
                bbox = block['bbox']
                # We'll extract the actual image data later if needed
                image_element = ImageElement(
                    bbox=tuple(bbox),
                    image_data=np.array([]),  # Placeholder
                    image_type='unknown'
                )
                image_elements.append(image_element)
        
        # Detect structural elements
        headings = self._detect_headings(text_elements)
        paragraphs = self._group_into_paragraphs(text_elements)
        bullet_points = self._detect_bullet_points(text_elements)
        table_candidates = self._detect_tables(text_elements)
        
        return {
            'page_number': page_number,
            'page_size': (page_width, page_height),
            'text_elements': text_elements,
            'image_elements': image_elements,
            'shape_elements': shapes,
            'table_elements': table_candidates,
            'headings': headings,
            'paragraphs': paragraphs,
            'bullet_points': bullet_points
        }

    def _extract_with_ocr(self, pdf_path: str, page_number: int) -> List[TextElement]:
        """
        Extract text using OCR for a specific page
        """
        try:
            # Convert specific page to image
            pages = convert_from_path(pdf_path, dpi=self.dpi, first_page=page_number, last_page=page_number)
            if not pages:
                return []
            
            page_image = pages[0]
            page_array = np.array(page_image)
            
            # Preprocess image for better OCR
            processed_image = self._preprocess_image(page_array)
            
            # Extract text using selected OCR engine
            if self.ocr_engine == 'easyocr' and EASYOCR_AVAILABLE:
                text_data = self._extract_with_easyocr_detailed(processed_image)
            else:
                text_data = self._extract_with_tesseract_detailed(processed_image)
            
            # Convert to TextElement objects
            text_elements = []
            for block in text_data:
                text_element = TextElement(
                    text=block['text'],
                    bbox=block['bbox'],
                    font_name='Arial',  # Default font for OCR
                    font_size=self._estimate_font_size(block['bbox']),
                    color=(0, 0, 0),  # Default black
                    bold=False,
                    italic=False,
                    confidence=block['confidence'] / 100.0,
                    alignment='left'
                )
                text_elements.append(text_element)
            
            return text_elements
            
        except Exception as e:
            print(f"ERROR: OCR extraction failed: {e}")
            return []

    def _extract_with_ocr_only(self, pdf_path: str) -> List[Dict]:
        """
        Fallback to OCR-only extraction
        """
        try:
            print(f"INFO: Using OCR-only extraction for: {pdf_path}")
            pages = convert_from_path(pdf_path, dpi=self.dpi)
            pages_data = []
            
            for page_num, page_image in enumerate(pages, 1):
                print(f"INFO: Processing page {page_num} with OCR...")
                
                page_array = np.array(page_image)
                processed_image = self._preprocess_image(page_array)
                
                # Extract text using selected OCR engine
                if self.ocr_engine == 'easyocr' and EASYOCR_AVAILABLE:
                    text_data = self._extract_with_easyocr_detailed(processed_image)
                else:
                    text_data = self._extract_with_tesseract_detailed(processed_image)
                
                # Convert to TextElement objects
                text_elements = []
                for block in text_data:
                    text_element = TextElement(
                        text=block['text'],
                        bbox=block['bbox'],
                        font_name='Arial',
                        font_size=self._estimate_font_size(block['bbox']),
                        color=(0, 0, 0),
                        bold=False,
                        italic=False,
                        confidence=block['confidence'] / 100.0,
                        alignment='left'
                    )
                    text_elements.append(text_element)
                
                page_data = {
                    'page_number': page_num,
                    'page_size': page_image.size,
                    'text_elements': text_elements,
                    'image_elements': [],
                    'shape_elements': [],
                    'table_elements': [],
                    'headings': self._detect_headings(text_elements),
                    'paragraphs': self._group_into_paragraphs(text_elements),
                    'bullet_points': self._detect_bullet_points(text_elements)
                }
                
                pages_data.append(page_data)
            
            return pages_data
            
        except Exception as e:
            print(f"ERROR: OCR-only extraction failed: {e}")
            traceback.print_exc()
            raise

    def _preprocess_image(self, image: np.ndarray) -> np.ndarray:
        """Enhanced image preprocessing for better OCR results"""
        # Convert to grayscale
        if len(image.shape) == 3:
            gray = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
        else:
            gray = image
        
        # Apply denoising
        denoised = cv2.fastNlMeansDenoising(gray)
        
        # Enhance contrast
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
        enhanced = clahe.apply(denoised)
        
        # Apply adaptive thresholding
        thresh = cv2.adaptiveThreshold(
            enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
        )
        
        return thresh

    def _extract_with_tesseract_detailed(self, image: np.ndarray) -> List[Dict]:
        """Extract text using Tesseract OCR with detailed information"""
        # Get detailed data from Tesseract
        data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT, config='--psm 6')
        
        text_blocks = []
        
        for i in range(len(data['text'])):
            text = data['text'][i].strip()
            conf = int(data['conf'][i]) if data['conf'][i] != -1 else 0
            
            if text and conf > 30:  # Filter low confidence text
                x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
                
                text_blocks.append({
                    'text': text,
                    'bbox': (x, y, x + w, y + h),
                    'confidence': conf
                })
        
        return text_blocks

    def _extract_with_easyocr_detailed(self, image: np.ndarray) -> List[Dict]:
        """Extract text using EasyOCR with detailed information"""
        results = self.reader.readtext(image)
        
        text_blocks = []
        for (bbox, text, confidence) in results:
            if confidence > 0.3:  # Filter low confidence text
                # Convert bbox format
                x1, y1 = int(min([point[0] for point in bbox])), int(min([point[1] for point in bbox]))
                x2, y2 = int(max([point[0] for point in bbox])), int(max([point[1] for point in bbox]))
                
                text_blocks.append({
                    'text': text,
                    'bbox': (x1, y1, x2, y2),
                    'confidence': int(confidence * 100)
                })
        
        return text_blocks

    def _int_to_rgb(self, color_int: int) -> Tuple[int, int, int]:
        """Convert integer color to RGB tuple"""
        if color_int == 0:
            return (0, 0, 0)  # Black
        
        # Extract RGB components
        r = (color_int >> 16) & 0xFF
        g = (color_int >> 8) & 0xFF
        b = color_int & 0xFF
        
        return (r, g, b)

    def _clean_font_name(self, font_name: str) -> str:
        """Clean and standardize font names"""
        if not font_name:
            return 'Arial'
        
        # Remove font subset prefixes (e.g., 'ABCDEF+Arial' -> 'Arial')
        if '+' in font_name:
            font_name = font_name.split('+')[-1]
        
        # Common font mappings
        font_mappings = {
            'TimesNewRomanPSMT': 'Times New Roman',
            'Times-Roman': 'Times New Roman',
            'Times-Bold': 'Times New Roman',
            'Times-Italic': 'Times New Roman',
            'ArialMT': 'Arial',
            'Arial-Bold': 'Arial',
            'Arial-Italic': 'Arial',
            'Helvetica': 'Arial',
            'Helvetica-Bold': 'Arial',
            'CourierNewPSMT': 'Courier New',
            'Courier': 'Courier New'
        }
        
        return font_mappings.get(font_name, font_name)

    def _determine_alignment_advanced(self, bbox: Tuple[float, float, float, float], page_width: float, text: str) -> str:
        """Advanced text alignment detection"""
        x1, _, x2, _ = bbox
        text_width = x2 - x1
        center_x = (x1 + x2) / 2
        
        # Check for center alignment (text centered on page)
        page_center = page_width / 2
        if abs(center_x - page_center) < page_width * 0.1:
            return 'center'
        
        # Check for right alignment (text near right margin)
        if x2 > page_width * 0.85:
            return 'right'
        
        # Check for left alignment (text near left margin)
        if x1 < page_width * 0.15:
            return 'left'
        
        # Default to left for body text
        return 'left'

    def _determine_alignment(self, bbox: Tuple[float, float, float, float], page_width: float) -> str:
        """Determine text alignment based on position"""
        x1, _, x2, _ = bbox
        center_x = (x1 + x2) / 2
        
        if center_x < page_width * 0.3:
            return 'left'
        elif center_x > page_width * 0.7:
            return 'right'
        else:
            return 'center'

    def _estimate_font_size(self, bbox: Tuple[int, int, int, int]) -> float:
        """Estimate font size from bounding box height"""
        _, y1, _, y2 = bbox
        height = y2 - y1
        # Rough estimation: font size is approximately 75% of height
        return max(8, height * 0.75)

    def _detect_headings(self, text_elements: List[TextElement]) -> List[TextElement]:
        """Detect heading elements based on font size and formatting"""
        if not text_elements:
            return []
        
        # Calculate average font size
        avg_font_size = sum(elem.font_size for elem in text_elements) / len(text_elements)
        
        headings = []
        for elem in text_elements:
            # Consider as heading if significantly larger than average or bold
            if elem.font_size > avg_font_size * 1.3 or elem.bold:
                headings.append(elem)
        
        return headings

    def _group_into_paragraphs(self, text_elements: List[TextElement]) -> List[List[TextElement]]:
        """Group text elements into paragraphs"""
        if not text_elements:
            return []
        
        # Sort by vertical position
        sorted_elements = sorted(text_elements, key=lambda x: x.bbox[1])
        
        paragraphs = []
        current_paragraph = []
        last_y = None
        
        for elem in sorted_elements:
            current_y = elem.bbox[1]
            
            # Start new paragraph if there's a significant vertical gap
            if last_y is not None and current_y - last_y > elem.font_size * 1.5:
                if current_paragraph:
                    paragraphs.append(current_paragraph)
                    current_paragraph = []
            
            current_paragraph.append(elem)
            last_y = current_y
        
        # Add the last paragraph
        if current_paragraph:
            paragraphs.append(current_paragraph)
        
        return paragraphs

    def _detect_bullet_points(self, text_elements: List[TextElement]) -> List[TextElement]:
        """Detect bullet point elements"""
        bullet_patterns = [r'^[•\-\*\+]\s+', r'^\d+[\.\)]\s+', r'^[a-zA-Z][\.\)]\s+']
        bullet_points = []
        
        for elem in text_elements:
            for pattern in bullet_patterns:
                if re.match(pattern, elem.text):
                    bullet_points.append(elem)
                    break
        
        return bullet_points

    def _detect_tables(self, text_elements: List[TextElement]) -> List[TableElement]:
        """Detect table structures"""
        # This is a simplified table detection
        # Group elements by similar Y coordinates (rows)
        rows = {}
        for elem in text_elements:
            y = int(elem.bbox[1] / 10) * 10  # Round to nearest 10 for grouping
            if y not in rows:
                rows[y] = []
            rows[y].append(elem)
        
        # Find rows with multiple columns (potential table rows)
        table_rows = []
        for y, elements in rows.items():
            if len(elements) >= 3:  # At least 3 columns
                # Sort by X position
                sorted_elements = sorted(elements, key=lambda x: x.bbox[0])
                table_rows.append(sorted_elements)
        
        # If we found potential table rows, create a table
        if len(table_rows) >= 2:
            # Combine all elements into one table
            all_elements = [elem for row in table_rows for elem in row]
            min_x = min(elem.bbox[0] for elem in all_elements)
            min_y = min(elem.bbox[1] for elem in all_elements)
            max_x = max(elem.bbox[2] for elem in all_elements)
            max_y = max(elem.bbox[3] for elem in all_elements)
            
            # Create table structure
            table_data = []
            cell_positions = []
            
            for row_elements in table_rows:
                row_texts = [elem.text for elem in row_elements]
                row_positions = [elem.bbox for elem in row_elements]
                table_data.append(row_texts)
                cell_positions.append(row_positions)
            
            table = TableElement(
                bbox=(min_x, min_y, max_x, max_y),
                rows=table_data,
                cell_positions=cell_positions
            )
            
            return [table]
        
        return []

    def create_powerpoint_presentation(self, pages_data: List[Dict], output_path: str, use_fallback_for_complex: bool = True):
        """
        Create PowerPoint presentation from extracted data
        """
        print(f"INFO: Creating enhanced PowerPoint presentation: {output_path}")
        
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Converted PDF Presentation"
        subtitle.text = f"Enhanced conversion with layout preservation • {len(pages_data)} pages"
        
        for page_data in pages_data:
            try:
                self._create_slide_from_page(prs, page_data, use_fallback_for_complex)
            except Exception as e:
                print(f"WARNING: Error creating slide for page {page_data['page_number']}: {e}")
                # Create a simple fallback slide
                self._create_fallback_slide(prs, page_data)
        
        prs.save(output_path)
        print(f"SUCCESS: Enhanced PowerPoint presentation saved: {output_path}")

    def _create_slide_from_page(self, prs: Presentation, page_data: Dict, use_fallback_for_complex: bool):
        """
        Create a slide from page data with exact formatting preservation
        """
        page_number = page_data['page_number']
        text_elements = page_data['text_elements']
        image_elements = page_data.get('image_elements', [])
        tables = page_data['table_elements']
        page_size = page_data['page_size']
        
        # Use blank layout for maximum control
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        print(f"INFO: Creating slide {page_number} with {len(text_elements)} text elements and {len(image_elements)} images")
        
        # First, add images with exact positioning
        for image_elem in image_elements:
            try:
                self._add_image_to_slide(slide, image_elem, page_size)
            except Exception as e:
                print(f"WARNING: Could not add image to slide {page_number}: {e}")
        
        # Add tables with exact positioning
        for table in tables:
            try:
                self._add_table_to_slide_exact(slide, table, page_size)
            except Exception as e:
                print(f"WARNING: Could not add table to slide {page_number}: {e}")
        
        # Add all text elements with exact positioning and formatting
        added_elements = 0
        for elem in text_elements:
            try:
                self._add_text_element_exact(slide, elem, page_size)
                added_elements += 1
            except Exception as e:
                print(f"WARNING: Could not add text element '{elem.text[:50]}' to slide {page_number}: {e}")
        
        print(f"INFO: Successfully added {added_elements}/{len(text_elements)} text elements to slide {page_number}")
        
        # If very few elements were added or slide is too complex, add fallback
        if use_fallback_for_complex and (added_elements < len(text_elements) * 0.5 or len(text_elements) > 100):
            print(f"INFO: Adding fallback background for complex page {page_number}")
            self._add_fallback_background_with_overlay(slide, page_data)

    def _add_table_to_slide(self, slide, table: TableElement):
        """Add a table to the slide"""
        try:
            rows = len(table.rows)
            cols = max(len(row) for row in table.rows) if table.rows else 1
            
            # Position table in the middle area of the slide
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            
            # Add table shape
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            ppt_table = table_shape.table
            
            # Fill table with data
            for i, row_data in enumerate(table.rows):
                for j, cell_text in enumerate(row_data):
                    if i < rows and j < cols:
                        cell = ppt_table.cell(i, j)
                        cell.text = cell_text
                        cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        except Exception as e:
            print(f"WARNING: Could not create table: {e}")

    def _add_image_to_slide(self, slide, image_elem: ImageElement, page_size: Tuple[float, float]):
        """Add image to slide with exact positioning"""
        try:
            # Convert PDF coordinates to PowerPoint coordinates
            pdf_width, pdf_height = page_size
            slide_width = Inches(10)  # Standard slide width
            slide_height = Inches(7.5)  # Standard slide height
            
            # Calculate scaling factors
            scale_x = slide_width.inches / pdf_width
            scale_y = slide_height.inches / pdf_height
            
            # Convert bbox coordinates
            x1, y1, x2, y2 = image_elem.bbox
            left = Inches(x1 * scale_x)
            top = Inches(y1 * scale_y)
            width = Inches((x2 - x1) * scale_x)
            height = Inches((y2 - y1) * scale_y)
            
            # Ensure coordinates are within slide bounds
            left = max(Inches(0), min(left, slide_width - Inches(0.1)))
            top = max(Inches(0), min(top, slide_height - Inches(0.1)))
            width = min(width, slide_width - left)
            height = min(height, slide_height - top)
            
            # Convert numpy array to PIL Image and save temporarily
            if image_elem.image_data.size > 0:
                if len(image_elem.image_data.shape) == 3:
                    image_pil = Image.fromarray(image_elem.image_data)
                else:
                    image_pil = Image.fromarray(image_elem.image_data, mode='L')
                
                # Save to temporary file
                import tempfile
                with tempfile.NamedTemporaryFile(suffix=f'.{image_elem.image_format}', delete=False) as tmp_file:
                    image_pil.save(tmp_file.name)
                    
                    # Add image to slide
                    pic = slide.shapes.add_picture(tmp_file.name, left, top, width, height)
                    print(f"INFO: Added image at ({left.inches:.2f}, {top.inches:.2f}) size ({width.inches:.2f}x{height.inches:.2f})")
                    
                    # Clean up temporary file
                    import os
                    os.unlink(tmp_file.name)
                    
        except Exception as e:
            print(f"WARNING: Could not add image: {e}")

    def _add_text_element_exact(self, slide, elem: TextElement, page_size: Tuple[float, float]):
        """Add single text element with exact positioning and formatting"""
        try:
            # Convert PDF coordinates to PowerPoint coordinates
            pdf_width, pdf_height = page_size
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            
            # Calculate scaling factors
            scale_x = slide_width.inches / pdf_width
            scale_y = slide_height.inches / pdf_height
            
            # Convert bbox coordinates
            x1, y1, x2, y2 = elem.bbox
            left = Inches(x1 * scale_x)
            top = Inches(y1 * scale_y)
            width = Inches(max(0.5, (x2 - x1) * scale_x))  # Minimum width
            height = Inches(max(0.2, (y2 - y1) * scale_y))  # Minimum height
            
            # Ensure coordinates are within slide bounds
            left = max(Inches(0), min(left, slide_width - Inches(0.1)))
            top = max(Inches(0), min(top, slide_height - Inches(0.1)))
            width = min(width, slide_width - left)
            height = min(height, slide_height - top)
            
            # Create text box
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = elem.text
            text_frame.word_wrap = True
            text_frame.auto_size = None
            
            # Apply exact formatting
            paragraph = text_frame.paragraphs[0]
            font = paragraph.font
            
            # Set font properties
            font.name = self._clean_font_name(elem.font_name)
            font.size = Pt(max(6, min(72, elem.font_size)))  # Reasonable bounds
            font.bold = elem.bold
            font.italic = elem.italic
            
            # Set color
            try:
                font.color.rgb = RGBColor(*elem.color)
            except:
                font.color.rgb = RGBColor(0, 0, 0)  # Default to black
            
            # Set alignment
            if elem.alignment == 'center':
                paragraph.alignment = PP_ALIGN.CENTER
            elif elem.alignment == 'right':
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.LEFT
            
            # Remove margins for exact positioning
            text_frame.margin_left = Inches(0)
            text_frame.margin_right = Inches(0)
            text_frame.margin_top = Inches(0)
            text_frame.margin_bottom = Inches(0)
            
        except Exception as e:
            raise Exception(f"Failed to add text element: {e}")

    def _add_table_to_slide_exact(self, slide, table: TableElement, page_size: Tuple[float, float]):
        """Add table to slide with exact positioning"""
        try:
            if not table.rows or not table.rows[0]:
                return
            
            # Convert PDF coordinates to PowerPoint coordinates
            pdf_width, pdf_height = page_size
            scale_x = Inches(10).inches / pdf_width
            scale_y = Inches(7.5).inches / pdf_height
            
            x1, y1, x2, y2 = table.bbox
            left = Inches(x1 * scale_x)
            top = Inches(y1 * scale_y)
            width = Inches((x2 - x1) * scale_x)
            height = Inches((y2 - y1) * scale_y)
            
            # Ensure reasonable bounds
            left = max(Inches(0.1), min(left, Inches(9)))
            top = max(Inches(0.1), min(top, Inches(6.5)))
            width = max(Inches(1), min(width, Inches(10) - left))
            height = max(Inches(0.5), min(height, Inches(7.5) - top))
            
            rows = len(table.rows)
            cols = max(len(row) for row in table.rows)
            
            # Add table shape
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            ppt_table = table_shape.table
            
            # Fill table with data and formatting
            for i, row_data in enumerate(table.rows):
                for j, cell_text in enumerate(row_data):
                    if i < rows and j < cols:
                        cell = ppt_table.cell(i, j)
                        cell.text = str(cell_text)
                        
                        # Format cell text
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(9)
                            paragraph.font.name = 'Arial'
            
            print(f"INFO: Added table with {rows}x{cols} cells at ({left.inches:.2f}, {top.inches:.2f})")
            
        except Exception as e:
            print(f"WARNING: Could not create exact table: {e}")

    def _create_page_background_image(self, pdf_path: str, page_number: int) -> Optional[str]:
        """Convert PDF page to background image"""
        try:
            # Convert PDF page to high-quality image
            pages = convert_from_path(
                pdf_path, 
                dpi=300,  # High DPI for quality
                first_page=page_number, 
                last_page=page_number
            )
            
            if not pages:
                return None
            
            page_image = pages[0]
            
            # Save to temporary file
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                page_image.save(tmp_file.name, 'PNG', quality=95)
                return tmp_file.name
                
        except Exception as e:
            print(f"WARNING: Could not create background image for page {page_number}: {e}")
            return None

    def _add_fallback_background_with_overlay(self, slide, page_data: Dict):
        """Add page as background image with text overlay for complex layouts"""
        try:
            page_number = page_data['page_number']
            
            # Try to create background image
            # This would require access to the original PDF path
            # For now, we'll create a comprehensive text overlay
            
            # Clear existing shapes (keep only essential ones)
            shapes_to_keep = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                    # Keep shapes with meaningful text
                    if len(shape.text_frame.text) > 5:
                        shapes_to_keep.append(shape)
            
            # Add comprehensive text content as structured overlay
            content_y = Inches(0.5)
            line_height = Inches(0.3)
            
            # Group text elements by approximate Y position for better layout
            text_elements = page_data.get('text_elements', [])
            if text_elements:
                # Sort by Y position
                sorted_elements = sorted(text_elements, key=lambda x: x.bbox[1])
                
                # Group into lines
                lines = []
                current_line = []
                last_y = None
                
                for elem in sorted_elements:
                    current_y = elem.bbox[1]
                    
                    # If significant Y difference, start new line
                    if last_y is not None and abs(current_y - last_y) > elem.font_size:
                        if current_line:
                            lines.append(current_line)
                            current_line = []
                    
                    current_line.append(elem)
                    last_y = current_y
                
                if current_line:
                    lines.append(current_line)
                
                # Create text boxes for each line with preserved formatting
                for line_elements in lines[:20]:  # Limit to first 20 lines to avoid overcrowding
                    if content_y.inches > 6.5:  # Stop if we run out of space
                        break
                    
                    # Combine elements in the line
                    line_text = ' '.join(elem.text for elem in line_elements)
                    if line_text.strip():
                        # Use formatting from the first element in the line
                        first_elem = line_elements[0]
                        
                        text_box = slide.shapes.add_textbox(
                            Inches(0.3), content_y, Inches(9.4), line_height
                        )
                        text_frame = text_box.text_frame
                        text_frame.text = line_text
                        text_frame.word_wrap = True
                        
                        # Apply formatting from original element
                        paragraph = text_frame.paragraphs[0]
                        font = paragraph.font
                        font.name = self._clean_font_name(first_elem.font_name)
                        font.size = Pt(max(8, min(14, first_elem.font_size * 0.8)))  # Slightly smaller
                        font.bold = first_elem.bold
                        font.italic = first_elem.italic
                        
                        try:
                            font.color.rgb = RGBColor(*first_elem.color)
                        except:
                            font.color.rgb = RGBColor(0, 0, 0)
                        
                        # Set alignment
                        if first_elem.alignment == 'center':
                            paragraph.alignment = PP_ALIGN.CENTER
                        elif first_elem.alignment == 'right':
                            paragraph.alignment = PP_ALIGN.RIGHT
                        else:
                            paragraph.alignment = PP_ALIGN.LEFT
                        
                        content_y = Inches(content_y.inches + line_height.inches)
            
            # Add page info at bottom
            info_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(7), Inches(9), Inches(0.4)
            )
            info_frame = info_box.text_frame
            info_frame.text = f"Page {page_number} - Layout preserved with {len(text_elements)} text elements"
            info_frame.paragraphs[0].font.size = Pt(8)
            info_frame.paragraphs[0].font.italic = True
            info_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
            
        except Exception as e:
            print(f"WARNING: Could not add fallback background: {e}")

    def _add_bullet_points_to_slide(self, slide, bullet_points: List[TextElement]):
        """Add bullet points to the slide"""
        if not bullet_points:
            return
        
        # Create text box for bullet points
        bullet_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(5)
        )
        text_frame = bullet_box.text_frame
        text_frame.word_wrap = True
        
        # Sort bullet points by vertical position
        sorted_bullets = sorted(bullet_points, key=lambda x: x.bbox[1])
        
        for i, bullet in enumerate(sorted_bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Clean bullet text (remove bullet markers)
            clean_text = re.sub(r'^[•\-\*\+\d+\w+][\.\)]\s*', '', bullet.text)
            p.text = clean_text
            p.font.size = Pt(max(12, bullet.font_size))
            p.font.color.rgb = RGBColor(*bullet.color)
            p.level = 0

    def _add_text_content_to_slide(self, slide, text_elements: List[TextElement], page_size: Tuple[float, float]):
        """Add text content preserving layout"""
        if not text_elements:
            return
        
        page_width, page_height = page_size
        
        # Group text elements by approximate regions
        for elem in text_elements:
            try:
                # Convert PDF coordinates to PowerPoint coordinates
                x1, y1, x2, y2 = elem.bbox
                
                # Convert to inches (rough conversion)
                left = Inches(x1 / page_width * 10)
                top = Inches(y1 / page_height * 7.5)
                width = Inches(max(1, (x2 - x1) / page_width * 10))
                height = Inches(max(0.3, (y2 - y1) / page_height * 7.5))
                
                # Ensure coordinates are within slide bounds
                left = max(Inches(0.1), min(left, Inches(9.5)))
                top = max(Inches(0.1), min(top, Inches(7)))
                width = min(width, Inches(10) - left)
                height = min(height, Inches(7.5) - top)
                
                # Create text box
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.text = elem.text
                text_frame.word_wrap = True
                
                # Apply formatting
                p = text_frame.paragraphs[0]
                p.font.size = Pt(max(8, min(20, elem.font_size)))
                p.font.color.rgb = RGBColor(*elem.color)
                p.font.bold = elem.bold
                p.font.italic = elem.italic
                
                # Set alignment
                if elem.alignment == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif elem.alignment == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT
                    
            except Exception as e:
                print(f"WARNING: Could not add text element '{elem.text[:50]}': {e}")

    def _add_fallback_content(self, slide, page_data: Dict):
        """Add fallback content when layout preservation fails"""
        # Create a simple text box with all content
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(6)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Add all text content as paragraphs
        all_text = []
        for paragraph in page_data['paragraphs']:
            paragraph_text = ' '.join(elem.text for elem in paragraph)
            all_text.append(paragraph_text)
        
        if all_text:
            content_frame.text = '\n\n'.join(all_text)
            content_frame.paragraphs[0].font.size = Pt(12)
        else:
            content_frame.text = f"Content from page {page_data['page_number']} - Complex layout preserved as text"

    def _create_fallback_slide(self, prs: Presentation, page_data: Dict):
        """Create a simple fallback slide"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = f"Page {page_data['page_number']}"
        
        # Combine all text
        all_text = []
        for elem in page_data['text_elements']:
            all_text.append(elem.text)
        
        content.text = '\n'.join(all_text) if all_text else "No text content found"

    def convert_pdf_to_powerpoint(self, input_pdf: str, output_pptx: str) -> str:
        """
        Main conversion function
        
        Args:
            input_pdf: Path to input PDF file
            output_pptx: Path for output PowerPoint file
            
        Returns:
            Path to the created PowerPoint file
        """
        try:
            print(f"INFO: Starting enhanced PDF to PowerPoint conversion")
            print(f"INFO: Input: {input_pdf}")
            print(f"INFO: Output: {output_pptx}")
            
            # Step 1: Detect PDF type
            pdf_type = self.detect_pdf_type(input_pdf)
            print(f"INFO: Detected PDF type: {pdf_type}")
            
            # Step 2: Extract content with formatting
            pages_data = self.extract_text_elements_with_formatting(input_pdf)
            print(f"INFO: Extracted content from {len(pages_data)} pages")
            
            # Step 2.5: Extract images
            page_images = self.extract_images_from_pdf(input_pdf)
            print(f"INFO: Extracted images from {len(page_images)} pages")
            
            # Merge images into pages_data
            for page_data in pages_data:
                page_num = page_data['page_number']
                if page_num in page_images:
                    page_data['image_elements'] = page_images[page_num]
                    print(f"INFO: Added {len(page_images[page_num])} images to page {page_num}")
            
            # Step 3: Create PowerPoint presentation
            self.create_powerpoint_presentation(pages_data, output_pptx)
            
            print(f"SUCCESS: Enhanced PDF to PowerPoint conversion completed")
            return output_pptx
            
        except Exception as e:
            print(f"ERROR: Enhanced conversion failed: {e}")
            traceback.print_exc()
            raise


def main():
    parser = argparse.ArgumentParser(description='Enhanced PDF to PowerPoint converter')
    parser.add_argument('input_pdf', help='Input PDF file path')
    parser.add_argument('output_pptx', help='Output PowerPoint file path')
    parser.add_argument('--ocr-engine', choices=['tesseract', 'easyocr'], default='tesseract',
                       help='OCR engine to use')
    parser.add_argument('--dpi', type=int, default=200, help='DPI for PDF to image conversion')
    
    args = parser.parse_args()
    
    try:
        converter = EnhancedPDFToPPTConverter(
            ocr_engine=args.ocr_engine,
            dpi=args.dpi
        )
        
        output_path = converter.convert_pdf_to_powerpoint(args.input_pdf, args.output_pptx)
        print(f"SUCCESS: Enhanced PowerPoint presentation created: {output_path}")
        
    except Exception as e:
        print(f"ERROR: Conversion failed: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
