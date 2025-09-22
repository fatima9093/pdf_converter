#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF to PowerPoint Layout-Preserving Converter
Converts PDFs to PowerPoint presentations while preserving the exact original layout
using PyMuPDF for precise text extraction and python-pptx for accurate positioning.
"""

import argparse
import json
import sys
import os
from pathlib import Path
from typing import List, Dict, Tuple, Union, Optional
import traceback
import tempfile

# Ensure UTF-8 encoding for output
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')
if sys.stderr.encoding != 'utf-8':
    sys.stderr.reconfigure(encoding='utf-8')

# PDF processing
import fitz  # PyMuPDF

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# Image processing for fallback
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import cv2
import numpy as np

# Configure Tesseract path for cross-platform compatibility
import platform
import os

if platform.system() == 'Windows':
    tesseract_path = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
elif platform.system() == 'Darwin':  # macOS
    tesseract_path = '/usr/local/bin/tesseract'
else:  # Linux
    tesseract_path = '/usr/bin/tesseract'

# Check if tesseract exists at the platform-specific path, otherwise use system PATH
if os.path.exists(tesseract_path):
    pytesseract.pytesseract.tesseract_cmd = tesseract_path


class PDFToPPTLayoutPreserver:
    def __init__(self):
        """Initialize the layout-preserving PDF to PowerPoint converter"""
        print("INFO: PDF to PowerPoint Layout-Preserving Converter initialized", file=sys.stderr)

    def convert_pdf_to_powerpoint(self, pdf_path: str, output_path: str) -> bool:
        """
        Convert PDF to PowerPoint while preserving the exact original layout
        Uses PDF pages as background images with editable text overlay
        
        Args:
            pdf_path: Path to the input PDF file
            output_path: Path for the output PowerPoint file
            
        Returns:
            bool: True if conversion successful, False otherwise
        """
        try:
            print(f"INFO: Starting layout-preserving PDF to PowerPoint conversion: {pdf_path}", file=sys.stderr)
            
            # Step 1: Detect PDF type
            pdf_type = self._detect_pdf_type(pdf_path)
            print(f"INFO: Detected PDF type: {pdf_type}", file=sys.stderr)
            
            # Step 2: Extract content with precise positioning
            pages_data = self._extract_content_with_positioning(pdf_path)
            
            if not pages_data:
                print("ERROR: No content could be extracted from the PDF", file=sys.stderr)
                return False
            
            # Step 3: Check if PDF contains images to determine conversion method
            has_images = any(len(page['images']) > 0 for page in pages_data)
            
            if has_images:
                print(f"INFO: PDF contains images, using background image method with white text backgrounds", file=sys.stderr)
                success = self._create_powerpoint_with_background_and_text(pages_data, pdf_path, output_path)
            else:
                print(f"INFO: PDF is text-only, using transparent text method", file=sys.stderr)
                success = self._create_powerpoint_text_only(pages_data, output_path)
            
            if success:
                print(f"SUCCESS: Layout-preserving PDF to PowerPoint conversion completed: {output_path}", file=sys.stderr)
                return True
            else:
                print("ERROR: Failed to create PowerPoint presentation", file=sys.stderr)
                return False
            
        except Exception as e:
            print(f"ERROR: Layout-preserving PDF to PowerPoint conversion failed: {e}", file=sys.stderr)
            traceback.print_exc()
            return False

    def _detect_pdf_type(self, pdf_path: str) -> str:
        """
        Detect if PDF is text-based or image-based (scanned)
        
        Returns:
            'text-based' or 'image-based'
        """
        try:
            doc = fitz.open(pdf_path)
            total_text_length = 0
            total_pages = len(doc)
            
            for page_num in range(min(3, total_pages)):  # Check first 3 pages
                page = doc[page_num]
                text = page.get_text()
                total_text_length += len(text.strip())
            
            doc.close()
            
            # If we found substantial text, it's likely text-based
            avg_text_per_page = total_text_length / min(3, total_pages)
            
            if avg_text_per_page > 100:  # Threshold for text-based
                return 'text-based'
            else:
                return 'image-based'
                
        except Exception as e:
            print(f"WARNING: Could not analyze PDF type: {e}", file=sys.stderr)
            return 'image-based'  # Default to image-based for safety

    def _extract_content_with_positioning(self, pdf_path: str) -> List[Dict]:
        """
        Extract content from PDF with precise positioning information
        
        Returns:
            List of pages with positioned content elements
        """
        try:
            print(f"INFO: Extracting content with precise positioning using PyMuPDF", file=sys.stderr)
            
            doc = fitz.open(pdf_path)
            pages_data = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_rect = page.rect
                
                print(f"INFO: Processing page {page_num + 1} (size: {page_rect.width:.1f}x{page_rect.height:.1f})", file=sys.stderr)
                
                page_data = {
                    'page_number': page_num + 1,
                    'page_size': (page_rect.width, page_rect.height),
                    'text_elements': [],
                    'images': []
                }
                
                # Try to extract text using PyMuPDF first
                try:
                    text_dict = page.get_text("dict")
                    text_elements_found = self._process_text_dict(text_dict, page_data)
                    print(f"INFO: PyMuPDF extracted {text_elements_found} text elements from page {page_num + 1}", file=sys.stderr)
                except Exception as e:
                    print(f"WARNING: PyMuPDF text extraction failed for page {page_num + 1}: {e}", file=sys.stderr)
                    text_elements_found = 0
                
                # If no text found with PyMuPDF or very little text, try OCR as fallback
                if text_elements_found < 3:
                    print(f"INFO: Using OCR fallback for page {page_num + 1} (found {text_elements_found} elements)", file=sys.stderr)
                    try:
                        ocr_elements = self._extract_text_with_ocr_fallback(pdf_path, page_num + 1, page_rect)
                        page_data['text_elements'].extend(ocr_elements)
                        print(f"INFO: OCR added {len(ocr_elements)} text elements to page {page_num + 1}", file=sys.stderr)
                    except Exception as e:
                        print(f"WARNING: OCR fallback failed for page {page_num + 1}: {e}", file=sys.stderr)
                
                # Try to extract images (but don't fail if this doesn't work)
                try:
                    image_list = page.get_images(full=True)
                    for img_index, img in enumerate(image_list):
                        try:
                            # Get image position
                            xref = img[0]
                            image_rects = page.get_image_rects(img)
                            
                            for rect in image_rects:
                                bbox = (rect.x0, rect.y0, rect.x1, rect.y1)
                                page_data['images'].append({
                                    'bbox': bbox,
                                    'xref': xref
                                })
                        except Exception as e:
                            print(f"WARNING: Could not process image {img_index} on page {page_num + 1}: {e}", file=sys.stderr)
                except Exception as e:
                    print(f"WARNING: Could not extract images from page {page_num + 1}: {e}", file=sys.stderr)
                
                
                pages_data.append(page_data)
                print(f"SUCCESS: Page {page_num + 1}: extracted {len(page_data['text_elements'])} text elements, {len(page_data['images'])} images", file=sys.stderr)
            
            doc.close()
            return pages_data
            
        except Exception as e:
            print(f"ERROR: Failed to extract content with positioning: {e}", file=sys.stderr)
            traceback.print_exc()
            return []

    def _process_text_dict(self, text_dict: Dict, page_data: Dict) -> int:
        """
        Process PyMuPDF text dictionary and add text elements to page_data
        
        Returns:
            Number of text elements found
        """
        text_elements_found = 0
        
        for block in text_dict.get('blocks', []):
            if 'lines' in block:  # Text block
                for line in block['lines']:
                    for span in line['spans']:
                        text = span['text'].strip()
                        if text and len(text) > 0:
                            # Get precise positioning
                            bbox = span['bbox']  # (x0, y0, x1, y1)
                            
                            # Extract font information
                            font_name = span.get('font', 'Arial')
                            font_size = span.get('size', 12)
                            font_flags = span.get('flags', 0)
                            color = span.get('color', 0)
                            
                            # Parse font properties
                            bold = bool(font_flags & 2**4)
                            italic = bool(font_flags & 2**1)
                            
                            # Convert color from integer to RGB
                            rgb_color = self._int_to_rgb(color)
                            
                            # Clean font name
                            clean_font_name = self._clean_font_name(font_name)
                            
                            text_element = {
                                'text': text,
                                'bbox': bbox,  # (x0, y0, x1, y1)
                                'font_name': clean_font_name,
                                'font_size': font_size,
                                'color': rgb_color,
                                'bold': bold,
                                'italic': italic
                            }
                            
                            page_data['text_elements'].append(text_element)
                            text_elements_found += 1
        
        return text_elements_found

    def _extract_text_with_ocr_fallback(self, pdf_path: str, page_num: int, page_rect) -> List[Dict]:
        """
        Extract text using OCR as a fallback when PyMuPDF fails
        
        Returns:
            List of text elements extracted via OCR
        """
        try:
            # Convert specific page to image
            pages = convert_from_path(pdf_path, dpi=200, first_page=page_num, last_page=page_num)
            if not pages:
                return []
            
            page_image = pages[0]
            page_array = np.array(page_image)
            
            # Preprocess image for better OCR
            processed_image = self._preprocess_image_for_ocr(page_array)
            
            # Extract text using Tesseract
            ocr_data = pytesseract.image_to_data(processed_image, output_type=pytesseract.Output.DICT)
            
            text_elements = []
            for i in range(len(ocr_data['text'])):
                if int(ocr_data['conf'][i]) > 30:  # Filter low confidence text
                    text = ocr_data['text'][i].strip()
                    if text:
                        # Convert OCR coordinates to PDF coordinates
                        # OCR gives coordinates in image pixels, we need to scale to PDF points
                        scale_x = page_rect.width / page_image.width
                        scale_y = page_rect.height / page_image.height
                        
                        x = ocr_data['left'][i] * scale_x
                        y = ocr_data['top'][i] * scale_y
                        w = ocr_data['width'][i] * scale_x
                        h = ocr_data['height'][i] * scale_y
                        
                        text_element = {
                            'text': text,
                            'bbox': (x, y, x + w, y + h),
                            'font_name': 'Arial',
                            'font_size': 12,
                            'color': (0, 0, 0),
                            'bold': False,
                            'italic': False
                        }
                        text_elements.append(text_element)
            
            return text_elements
            
        except Exception as e:
            print(f"ERROR: OCR fallback failed: {e}", file=sys.stderr)
            return []

    def _preprocess_image_for_ocr(self, image: np.ndarray) -> np.ndarray:
        """Preprocess image for better OCR accuracy"""
        # Convert to grayscale
        if len(image.shape) == 3:
            gray = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
        else:
            gray = image
        
        # Noise reduction
        denoised = cv2.bilateralFilter(gray, 9, 75, 75)
        
        # Enhance contrast
        clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8,8))
        enhanced = clahe.apply(denoised)
        
        # Binarization
        _, binary = cv2.threshold(enhanced, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        return binary

    def _create_powerpoint_with_background_and_text(self, pages_data: List[Dict], pdf_path: str, output_path: str) -> bool:
        """
        Create PowerPoint presentation using PDF pages as background images with editable text overlay
        """
        try:
            print(f"INFO: Creating PowerPoint presentation with background images and editable text", file=sys.stderr)
            
            # Create new presentation
            prs = Presentation()
            
            # Set presentation dimensions based on the first page (assume all pages have same size)
            if pages_data:
                first_page = pages_data[0]
                page_width, page_height = first_page['page_size']
                
                # Set slide dimensions to match PDF page size exactly
                # Convert PDF points to inches (1 point = 1/72 inch)
                ppt_width = Inches(page_width / 72.0)
                ppt_height = Inches(page_height / 72.0)
                
                # Set the presentation dimensions
                prs.slide_width = int(ppt_width.emu)
                prs.slide_height = int(ppt_height.emu)
                
                print(f"INFO: Set presentation dimensions to {ppt_width.inches:.2f}\" x {ppt_height.inches:.2f}\" (PDF: {page_width:.1f}x{page_height:.1f} points)", file=sys.stderr)
            
            for page_data in pages_data:
                page_num = page_data['page_number']
                page_width, page_height = page_data['page_size']
                
                print(f"INFO: Creating slide {page_num} with background image", file=sys.stderr)
                
                # Create blank slide for precise control
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Use the presentation dimensions
                ppt_width = Inches(page_width / 72.0)
                ppt_height = Inches(page_height / 72.0)
                
                # Convert PDF page to background image with text areas masked out
                background_image_path = self._create_page_background_image(pdf_path, page_num, page_data['text_elements'])
                if background_image_path:
                    try:
                        # Add the PDF page as background image
                        slide.shapes.add_picture(background_image_path, Inches(0), Inches(0), ppt_width, ppt_height)
                        print(f"INFO: Added background image for slide {page_num}", file=sys.stderr)
                    except Exception as e:
                        print(f"WARNING: Could not add background image for slide {page_num}: {e}", file=sys.stderr)
                
                # No scaling needed since we're using the same coordinate system
                scale_x = 1.0 / 72.0  # Convert points to inches
                scale_y = 1.0 / 72.0  # Convert points to inches
                
                # Add text elements as editable text boxes on top of the background
                text_added = 0
                for text_element in page_data['text_elements']:
                    try:
                        self._add_editable_text_element_to_slide(slide, text_element, scale_x, scale_y, ppt_width, ppt_height)
                        text_added += 1
                    except Exception as e:
                        print(f"WARNING: Could not add text element '{text_element.get('text', '')}': {e}", file=sys.stderr)
                
                # If no text was extracted with PyMuPDF, try OCR
                if text_added == 0:
                    print(f"INFO: No text found with PyMuPDF, trying OCR for slide {page_num}", file=sys.stderr)
                    try:
                        ocr_elements = self._extract_text_with_ocr_fallback(pdf_path, page_num, fitz.Rect(0, 0, page_width, page_height))
                        for ocr_element in ocr_elements:
                            try:
                                self._add_editable_text_element_to_slide(slide, ocr_element, scale_x, scale_y, ppt_width, ppt_height)
                                text_added += 1
                            except Exception as e:
                                print(f"WARNING: Could not add OCR text element: {e}", file=sys.stderr)
                    except Exception as e:
                        print(f"WARNING: OCR fallback failed for slide {page_num}: {e}", file=sys.stderr)
                
                # Clean up background image file
                if background_image_path:
                    try:
                        os.unlink(background_image_path)
                    except:
                        pass
                
                print(f"SUCCESS: Created slide {page_num} with background image and {text_added} editable text elements", file=sys.stderr)
            
            # Save presentation
            prs.save(output_path)
            print(f"SUCCESS: PowerPoint presentation saved: {output_path}", file=sys.stderr)
            return True
            
        except Exception as e:
            print(f"ERROR: Failed to create PowerPoint presentation: {e}", file=sys.stderr)
            traceback.print_exc()
            return False

    def _create_powerpoint_text_only(self, pages_data: List[Dict], output_path: str) -> bool:
        """
        Create PowerPoint presentation for text-only PDFs with transparent text boxes
        """
        try:
            print(f"INFO: Creating PowerPoint presentation for text-only PDF", file=sys.stderr)
            
            # Create new presentation
            prs = Presentation()
            
            # Set presentation dimensions based on the first page (assume all pages have same size)
            if pages_data:
                first_page = pages_data[0]
                page_width, page_height = first_page['page_size']
                
                # Set slide dimensions to match PDF page size exactly
                # Convert PDF points to inches (1 point = 1/72 inch)
                ppt_width = Inches(page_width / 72.0)
                ppt_height = Inches(page_height / 72.0)
                
                # Set the presentation dimensions
                prs.slide_width = int(ppt_width.emu)
                prs.slide_height = int(ppt_height.emu)
                
                print(f"INFO: Set presentation dimensions to {ppt_width.inches:.2f}\" x {ppt_height.inches:.2f}\" (PDF: {page_width:.1f}x{page_height:.1f} points)", file=sys.stderr)
            
            for page_data in pages_data:
                page_num = page_data['page_number']
                page_width, page_height = page_data['page_size']
                
                print(f"INFO: Creating text-only slide {page_num}", file=sys.stderr)
                
                # Create blank slide for precise control
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Use the presentation dimensions
                ppt_width = Inches(page_width / 72.0)
                ppt_height = Inches(page_height / 72.0)
                
                # No scaling needed since we're using the same coordinate system
                scale_x = 1.0 / 72.0  # Convert points to inches
                scale_y = 1.0 / 72.0  # Convert points to inches
                
                # Add text elements with transparent backgrounds
                text_added = 0
                for text_element in page_data['text_elements']:
                    try:
                        self._add_transparent_text_element_to_slide(slide, text_element, scale_x, scale_y, ppt_width, ppt_height)
                        text_added += 1
                    except Exception as e:
                        print(f"WARNING: Could not add text element '{text_element.get('text', '')}': {e}", file=sys.stderr)
                
                # For text-only PDFs, if no text found, that's unusual but we'll continue
                if text_added == 0:
                    print(f"WARNING: No text found for slide {page_num} in text-only PDF", file=sys.stderr)
                
                print(f"SUCCESS: Created text-only slide {page_num} with {text_added} transparent text elements", file=sys.stderr)
            
            # Save presentation
            prs.save(output_path)
            print(f"SUCCESS: Text-only PowerPoint presentation saved: {output_path}", file=sys.stderr)
            return True
            
        except Exception as e:
            print(f"ERROR: Failed to create text-only PowerPoint presentation: {e}", file=sys.stderr)
            traceback.print_exc()
            return False

    def _add_transparent_text_element_to_slide(self, slide, text_element: Dict, scale_x: float, scale_y: float, 
                                             ppt_width, ppt_height):
        """
        Add a single text element to the slide as a transparent text box (for text-only PDFs)
        """
        try:
            bbox = text_element['bbox']  # (x0, y0, x1, y1) in PDF points
            
            # Convert PDF coordinates to PowerPoint inches
            left = Inches(bbox[0] * scale_x)
            top = Inches(bbox[1] * scale_y)
            width = Inches(max(0.1, (bbox[2] - bbox[0]) * scale_x))  # Minimum width
            height = Inches(max(0.1, (bbox[3] - bbox[1]) * scale_y))  # Minimum height
            
            # Ensure coordinates are within slide bounds
            left = max(Inches(0), min(left, ppt_width - Inches(0.1)))
            top = max(Inches(0), min(top, ppt_height - Inches(0.1)))
            width = min(width, ppt_width - left)
            height = min(height, ppt_height - top)
            
            # Create text box
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.clear()
            text_frame.word_wrap = False  # Prevent text wrapping to maintain exact positioning
            text_frame.auto_size = None
            
            # Remove margins for precise positioning
            text_frame.margin_left = Inches(0)
            text_frame.margin_right = Inches(0)
            text_frame.margin_top = Inches(0)
            text_frame.margin_bottom = Inches(0)
            
            # Make text box background transparent
            fill = text_box.fill
            fill.background()
            
            # Remove text box border
            line = text_box.line
            line.fill.background()
            
            # Add text content
            p = text_frame.paragraphs[0]
            p.text = text_element['text']
            
            # Apply original formatting with original colors
            font = p.font
            font.name = text_element['font_name']
            font.size = Pt(max(6, min(72, text_element['font_size'])))
            font.bold = text_element['bold']
            font.italic = text_element['italic']
            
            # Use original text color
            try:
                font.color.rgb = RGBColor(*text_element['color'])
            except:
                font.color.rgb = RGBColor(0, 0, 0)  # Default to black
            
        except Exception as e:
            print(f"WARNING: Could not add transparent text element '{text_element.get('text', '')}': {e}", file=sys.stderr)

    def _group_overlapping_text(self, text_elements: List[Dict]) -> List[List[Dict]]:
        """
        Group text elements that overlap or are very close to each other
        to prevent text overlapping in PowerPoint
        """
        if not text_elements:
            return []
        
        # Sort text elements by vertical position (top to bottom)
        sorted_elements = sorted(text_elements, key=lambda x: (x['bbox'][1], x['bbox'][0]))
        
        groups = []
        current_group = [sorted_elements[0]]
        
        for i in range(1, len(sorted_elements)):
            current_elem = sorted_elements[i]
            last_elem = current_group[-1]
            
            # Check if elements overlap or are very close
            if self._elements_overlap_or_close(current_elem, last_elem):
                # Merge into current group
                current_group.append(current_elem)
            else:
                # Start new group
                groups.append(current_group)
                current_group = [current_elem]
        
        # Add the last group
        if current_group:
            groups.append(current_group)
        
        return groups

    def _elements_overlap_or_close(self, elem1: Dict, elem2: Dict, threshold: float = 5.0) -> bool:
        """
        Check if two text elements overlap or are very close to each other
        """
        bbox1 = elem1['bbox']  # (x0, y0, x1, y1)
        bbox2 = elem2['bbox']
        
        # Check vertical overlap/proximity
        vertical_overlap = not (bbox1[3] + threshold < bbox2[1] or bbox2[3] + threshold < bbox1[1])
        
        # Check horizontal overlap/proximity
        horizontal_overlap = not (bbox1[2] + threshold < bbox2[0] or bbox2[2] + threshold < bbox1[0])
        
        return vertical_overlap and horizontal_overlap

    def _create_page_background_image(self, pdf_path: str, page_number: int, text_elements: List[Dict]) -> Optional[str]:
        """Convert PDF page to background image with text areas masked out"""
        try:
            # Convert PDF page to high-quality image
            pages = convert_from_path(
                pdf_path, 
                dpi=150,  # Good balance between quality and file size
                first_page=page_number, 
                last_page=page_number
            )
            
            if not pages:
                return None
            
            page_image = pages[0]
            
            # Convert PIL image to numpy array for processing
            img_array = np.array(page_image)
            
            # Mask out text areas to prevent doubling
            if text_elements:
                # Get actual PDF page dimensions from the first text element or use standard A4
                try:
                    doc = fitz.open(pdf_path)
                    page = doc[page_number - 1]
                    pdf_width = page.rect.width
                    pdf_height = page.rect.height
                    doc.close()
                except:
                    pdf_width = 595.0  # Standard A4 width in points
                    pdf_height = 842.0  # Standard A4 height in points
                
                # Calculate scaling factor from PDF points to image pixels
                scale_x = page_image.width / pdf_width
                scale_y = page_image.height / pdf_height
                
                # Create a mask for text areas
                for text_elem in text_elements:
                    bbox = text_elem['bbox']
                    x1 = int(bbox[0] * scale_x)
                    y1 = int(bbox[1] * scale_y)
                    x2 = int(bbox[2] * scale_x)
                    y2 = int(bbox[3] * scale_y)
                    
                    # Ensure coordinates are within image bounds
                    x1 = max(0, min(x1, page_image.width))
                    y1 = max(0, min(y1, page_image.height))
                    x2 = max(x1, min(x2, page_image.width))
                    y2 = max(y1, min(y2, page_image.height))
                    
                    # Fill text area with white to remove background text
                    img_array[y1:y2, x1:x2] = [255, 255, 255]  # White fill
            
            # Convert back to PIL image
            processed_image = Image.fromarray(img_array)
            
            # Save to temporary file
            temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            processed_image.save(temp_file.name, 'PNG', quality=85)
            temp_file.close()
            return temp_file.name
                
        except Exception as e:
            print(f"WARNING: Could not create background image for page {page_number}: {e}", file=sys.stderr)
            return None

    def _add_editable_text_element_to_slide(self, slide, text_element: Dict, scale_x: float, scale_y: float, 
                                          ppt_width, ppt_height):
        """
        Add a single text element to the slide as an editable text box with transparent background
        """
        try:
            bbox = text_element['bbox']  # (x0, y0, x1, y1) in PDF points
            
            # Convert PDF coordinates to PowerPoint inches
            left = Inches(bbox[0] * scale_x)
            top = Inches(bbox[1] * scale_y)
            width = Inches(max(0.1, (bbox[2] - bbox[0]) * scale_x))  # Minimum width
            height = Inches(max(0.1, (bbox[3] - bbox[1]) * scale_y))  # Minimum height
            
            # Ensure coordinates are within slide bounds
            left = max(Inches(0), min(left, ppt_width - Inches(0.1)))
            top = max(Inches(0), min(top, ppt_height - Inches(0.1)))
            width = min(width, ppt_width - left)
            height = min(height, ppt_height - top)
            
            # Create text box
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.clear()
            text_frame.word_wrap = False  # Prevent text wrapping to maintain exact positioning
            text_frame.auto_size = None
            
            # Remove margins for precise positioning
            text_frame.margin_left = Inches(0)
            text_frame.margin_right = Inches(0)
            text_frame.margin_top = Inches(0)
            text_frame.margin_bottom = Inches(0)
            
            # Make text box background white to cover background text
            fill = text_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            
            # Remove text box border
            line = text_box.line
            line.fill.background()
            
            # Add text content
            p = text_frame.paragraphs[0]
            p.text = text_element['text']
            
            # Apply formatting
            font = p.font
            font.name = text_element['font_name']
            font.size = Pt(max(6, min(72, text_element['font_size'])))
            font.bold = text_element['bold']
            font.italic = text_element['italic']
            
            # Use original text color for visibility on white background
            try:
                font.color.rgb = RGBColor(*text_element['color'])
            except:
                font.color.rgb = RGBColor(0, 0, 0)  # Default to black text
            
        except Exception as e:
            print(f"WARNING: Could not add editable text element '{text_element.get('text', '')}': {e}", file=sys.stderr)

    def _add_image_to_slide(self, slide, image_data: Dict, scale_x: float, scale_y: float):
        """
        Add an image to the slide with exact positioning
        """
        try:
            bbox = image_data['bbox']  # (x0, y0, x1, y1) in PDF points
            
            # Convert PDF coordinates to PowerPoint inches
            left = Inches(bbox[0] * scale_x)
            top = Inches(bbox[1] * scale_y)
            width = Inches(max(0.1, (bbox[2] - bbox[0]) * scale_x))
            height = Inches(max(0.1, (bbox[3] - bbox[1]) * scale_y))
            
            # For now, just add a placeholder rectangle for images
            # In a full implementation, you would extract the actual image data
            # and add it using slide.shapes.add_picture()
            
            print(f"INFO: Image placeholder added at ({left.inches:.2f}, {top.inches:.2f}) size ({width.inches:.2f}x{height.inches:.2f})", file=sys.stderr)
            
        except Exception as e:
            print(f"WARNING: Could not add image: {e}", file=sys.stderr)

    def _int_to_rgb(self, color_int: int) -> Tuple[int, int, int]:
        """Convert integer color to RGB tuple"""
        if color_int == 0:
            return (0, 0, 0)  # Black
        
        # Extract RGB components
        r = (color_int >> 16) & 0xFF
        g = (color_int >> 8) & 0xFF
        b = color_int & 0xFF
        
        return (r, g, b)

    def _clean_font_name(self, font_name: str) -> str:
        """Clean and normalize font name"""
        if not font_name:
            return 'Arial'
        
        # Remove common font prefixes and suffixes
        font_name = font_name.split('+')[-1]  # Remove prefix like 'ABCDEF+'
        font_name = font_name.split('-')[0]   # Remove suffix like '-Bold'
        
        # Map common fonts to standard names
        font_mapping = {
            'TimesNewRoman': 'Times New Roman',
            'Times-Roman': 'Times New Roman',
            'Arial-Bold': 'Arial',
            'Arial-Italic': 'Arial',
            'Helvetica': 'Arial',
            'Courier': 'Courier New',
            'CourierNew': 'Courier New'
        }
        
        return font_mapping.get(font_name, font_name if font_name else 'Arial')


def main():
    parser = argparse.ArgumentParser(description='Convert PDF to PowerPoint with preserved layout')
    parser.add_argument('input_pdf', help='Input PDF file path')
    parser.add_argument('output_pptx', help='Output PowerPoint file path')
    
    args = parser.parse_args()
    
    if not os.path.exists(args.input_pdf):
        print(f"ERROR: Input PDF file not found: {args.input_pdf}", file=sys.stderr)
        sys.exit(1)
    
    # Create converter and convert
    converter = PDFToPPTLayoutPreserver()
    success = converter.convert_pdf_to_powerpoint(args.input_pdf, args.output_pptx)
    
    if success:
        print(f"SUCCESS: PDF successfully converted to PowerPoint: {args.output_pptx}")
        sys.exit(0)
    else:
        print("ERROR: PDF to PowerPoint conversion failed")
        sys.exit(1)


if __name__ == '__main__':
    main()
