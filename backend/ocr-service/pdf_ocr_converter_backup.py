#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PDF OCR Conversion Service
Converts scanned/image-based PDFs to Office documents using OCR
"""

import argparse
import json
import sys
import os
from pathlib import Path
from typing import List, Dict, Tuple, Union
import traceback

# Ensure UTF-8 encoding for output
if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8')
if sys.stderr.encoding != 'utf-8':
    sys.stderr.reconfigure(encoding='utf-8')

# PDF and image processing
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import cv2
import numpy as np

# Configure Tesseract path for cross-platform compatibility
import platform
import os

if platform.system() == 'Windows':
    tesseract_path = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
elif platform.system() == 'Darwin':  # macOS
    tesseract_path = '/usr/local/bin/tesseract'
else:  # Linux
    tesseract_path = '/usr/bin/tesseract'

# Check if tesseract exists at the platform-specific path, otherwise use system PATH
if os.path.exists(tesseract_path):
    pytesseract.pytesseract.tesseract_cmd = tesseract_path

# Try to import EasyOCR, but don't fail if it's not available
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError as e:
    print(f"WARNING: EasyOCR not available: {e}")
    print("INFO: Will use Tesseract OCR only")
    easyocr = None
    EASYOCR_AVAILABLE = False

# Document generation
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import openpyxl
from openpyxl.styles import Font, Alignment
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN


class PDFOCRConverter:
    def __init__(self, ocr_engine='tesseract'):
        """
        Initialize the PDF OCR converter
        
        Args:
            ocr_engine: 'tesseract' or 'easyocr'
        """
        # Check if EasyOCR is requested but not available
        if ocr_engine == 'easyocr' and not EASYOCR_AVAILABLE:
            print("WARNING: EasyOCR requested but not available, falling back to Tesseract")
            ocr_engine = 'tesseract'
        
        self.ocr_engine = ocr_engine
        self.reader = None
        
        if ocr_engine == 'easyocr' and EASYOCR_AVAILABLE:
            print("INFO: Initializing EasyOCR...")
            self.reader = easyocr.Reader(['en'])
        else:
            print("INFO: Using Tesseract OCR")
            
        print(f"SUCCESS: OCR Converter initialized with {ocr_engine}")

    def extract_text_from_pdf(self, pdf_path: str, dpi: int = 200) -> List[Dict]:
        """
        Extract text from PDF using OCR
        
        Args:
            pdf_path: Path to the PDF file
            dpi: Resolution for PDF to image conversion
            
        Returns:
            List of pages with extracted text and layout info
        """
        try:
            print(f"INFO: Converting PDF to images: {pdf_path}")
            
            # Convert PDF pages to images
            pages = convert_from_path(pdf_path, dpi=dpi)
            print(f"INFO: Found {len(pages)} page(s)")
            
            extracted_pages = []
            
            for page_num, page_image in enumerate(pages, 1):
                print(f"INFO: Processing page {page_num}...")
                
                # Convert PIL image to numpy array for OpenCV
                page_array = np.array(page_image)
                
                # Preprocess image for better OCR
                processed_image = self._preprocess_image(page_array)
                
                # Extract text using selected OCR engine
                if self.ocr_engine == 'easyocr' and EASYOCR_AVAILABLE:
                    text_data = self._extract_with_easyocr(processed_image)
                else:
                    text_data = self._extract_with_tesseract(processed_image)
                
                page_info = {
                    'page_number': page_num,
                    'text_blocks': text_data,
                    'full_text': ' '.join([block['text'] for block in text_data if block['text'].strip()]),
                    'image_size': page_image.size
                }
                
                extracted_pages.append(page_info)
                print(f"SUCCESS: Page {page_num}: {len(text_data)} text blocks, {len(page_info['full_text'])} characters")
            
            return extracted_pages
            
        except Exception as e:
            print(f"ERROR: Error extracting text from PDF: {e}")
            traceback.print_exc()
            raise

    def _preprocess_image(self, image: np.ndarray) -> np.ndarray:
        """Preprocess image for better OCR results"""
        # Convert to grayscale
        if len(image.shape) == 3:
            gray = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
        else:
            gray = image
        
        # Apply denoising
        denoised = cv2.fastNlMeansDenoising(gray)
        
        # Apply adaptive thresholding
        thresh = cv2.adaptiveThreshold(
            denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
        )
        
        return thresh

    def _extract_with_tesseract(self, image: np.ndarray) -> List[Dict]:
        """Extract text using Tesseract OCR"""
        # Get detailed data from Tesseract
        data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
        
        text_blocks = []
        current_block = {'text': '', 'bbox': None, 'confidence': 0}
        
        for i in range(len(data['text'])):
            text = data['text'][i].strip()
            conf = int(data['conf'][i])
            
            if text and conf > 30:  # Filter low confidence text
                x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
                
                text_blocks.append({
                    'text': text,
                    'bbox': (x, y, x + w, y + h),
                    'confidence': conf
                })
        
        return text_blocks

    def _extract_with_easyocr(self, image: np.ndarray) -> List[Dict]:
        """Extract text using EasyOCR"""
        results = self.reader.readtext(image)
        
        text_blocks = []
        for (bbox, text, confidence) in results:
            if confidence > 0.3:  # Filter low confidence text
                # Convert bbox format
                x1, y1 = int(min([point[0] for point in bbox])), int(min([point[1] for point in bbox]))
                x2, y2 = int(max([point[0] for point in bbox])), int(max([point[1] for point in bbox]))
                
                text_blocks.append({
                    'text': text,
                    'bbox': (x1, y1, x2, y2),
                    'confidence': int(confidence * 100)
                })
        
        return text_blocks

    def create_word_document(self, pages_data: List[Dict], output_path: str):
        """Create Word document from OCR data"""
        print(f"INFO: Creating Word document: {output_path}")
        
        doc = Document()
        
        # Add title
        title = doc.add_heading('OCR Converted Document', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        for page_data in pages_data:
            # Add page break for pages after the first
            if page_data['page_number'] > 1:
                doc.add_page_break()
            
            # Add page header
            page_header = doc.add_heading(f'Page {page_data["page_number"]}', level=2)
            
            # Group text blocks into paragraphs based on vertical proximity
            paragraphs = self._group_text_into_paragraphs(page_data['text_blocks'])
            
            for paragraph_text in paragraphs:
                if paragraph_text.strip():
                    p = doc.add_paragraph(paragraph_text)
                    p.style.font.size = Pt(11)
        
        doc.save(output_path)
        print(f"SUCCESS: Word document saved: {output_path}")

    def create_excel_document(self, pages_data: List[Dict], output_path: str):
        """Create Excel document from OCR data"""
        print(f"INFO: Creating Excel document: {output_path}")
        
        wb = openpyxl.Workbook()
        
        for page_data in pages_data:
            # Create worksheet for each page
            ws_name = f"Page_{page_data['page_number']}"
            if page_data['page_number'] == 1:
                ws = wb.active
                ws.title = ws_name
            else:
                ws = wb.create_sheet(ws_name)
            
            # Try to organize text into table-like structure
            table_data = self._organize_text_for_excel(page_data['text_blocks'])
            
            # Write data to worksheet
            for row_idx, row_data in enumerate(table_data, 1):
                for col_idx, cell_value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                    cell.font = Font(size=10)
                    cell.alignment = Alignment(wrap_text=True, vertical='top')
            
            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
        
        wb.save(output_path)
        print(f"SUCCESS: Excel document saved: {output_path}")

    def create_powerpoint_document(self, pages_data: List[Dict], output_path: str):
        """Create PowerPoint document from OCR data"""
        print(f"INFO: Creating PowerPoint document: {output_path}")
        
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "OCR Converted Presentation"
        subtitle.text = f"Converted from PDF • {len(pages_data)} pages"
        
        for page_data in pages_data:
            # Add slide for each page
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"Page {page_data['page_number']}"
            title_frame.paragraphs[0].font.size = PptxPt(24)
            title_frame.paragraphs[0].font.bold = True
            
            # Add content
            content_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(1.5), PptxInches(9), PptxInches(6)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Group text into bullet points
            paragraphs = self._group_text_into_paragraphs(page_data['text_blocks'])
            
            for i, paragraph_text in enumerate(paragraphs):
                if paragraph_text.strip():
                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = paragraph_text
                    p.font.size = PptxPt(14)
                    p.level = 0
        
        prs.save(output_path)
        print(f"SUCCESS: PowerPoint document saved: {output_path}")

    def _group_text_into_paragraphs(self, text_blocks: List[Dict]) -> List[str]:
        """Group text blocks into logical paragraphs"""
        if not text_blocks:
            return []
        
        # Sort text blocks by vertical position (top to bottom)
        sorted_blocks = sorted(text_blocks, key=lambda x: x['bbox'][1])
        
        paragraphs = []
        current_paragraph = []
        last_y = None
        
        for block in sorted_blocks:
            text = block['text'].strip()
            if not text:
                continue
                
            bbox = block['bbox']
            current_y = bbox[1]  # Top coordinate
            
            # Start new paragraph if there's a significant vertical gap
            if last_y is not None and current_y - last_y > 20:
                if current_paragraph:
                    paragraphs.append(' '.join(current_paragraph))
                    current_paragraph = []
            
            current_paragraph.append(text)
            last_y = current_y
        
        # Add the last paragraph
        if current_paragraph:
            paragraphs.append(' '.join(current_paragraph))
        
        return paragraphs

    def _organize_text_for_excel(self, text_blocks: List[Dict]) -> List[List[str]]:
        """Organize text blocks into table-like structure for Excel with proper column alignment"""
        if not text_blocks:
            return [["No text found"]]
        
        # Extract text blocks with position information
        positioned_blocks = []
        for block in text_blocks:
            text = block['text'].strip()
            if text:
                positioned_blocks.append({
                    'text': text,
                    'left': block['bbox'][0],  # x coordinate
                    'top': block['bbox'][1],   # y coordinate
                    'right': block['bbox'][2],
                    'bottom': block['bbox'][3]
                })
        
        if not positioned_blocks:
            return [["No text found"]]
        
        # Step 1: Determine column positions by analyzing horizontal alignment
        column_positions = self._determine_column_positions(positioned_blocks)
        print(f"DEBUG: Detected {len(column_positions)} columns at positions: {column_positions}")
        
        # Step 2: Group text blocks into rows based on vertical position
        rows = self._group_blocks_into_rows(positioned_blocks)
        print(f"DEBUG: Grouped text into {len(rows)} rows")
        
        # Step 3: Create table structure with proper column alignment
        table = []
        for row_blocks in rows:
            row_data = [''] * len(column_positions)
            
            for block in row_blocks:
                col_index = self._find_best_column(block['left'], column_positions)
                if 0 <= col_index < len(row_data):
                    if row_data[col_index]:
                        row_data[col_index] += ' ' + block['text']
                    else:
                        row_data[col_index] = block['text']
            
            # Only add rows that have actual content
            if any(cell.strip() for cell in row_data):
                table.append(row_data)
        
        print(f"DEBUG: Created table with {len(table)} rows and {len(column_positions)} columns")
        return table if table else [["No text found"]]

    def _determine_column_positions(self, blocks: List[Dict]) -> List[int]:
        """Determine column positions based on text block horizontal alignment"""
        # Collect all left positions
        left_positions = [block['left'] for block in blocks]
        
        if not left_positions:
            return [0]
        
        # Sort and cluster similar positions
        left_positions.sort()
        columns = [left_positions[0]]
        tolerance = 25  # Pixel tolerance for same column
        
        for pos in left_positions[1:]:
            # Check if this position is far enough from existing columns
            if all(abs(pos - col) > tolerance for col in columns):
                columns.append(pos)
        
        return sorted(columns)

    def _group_blocks_into_rows(self, blocks: List[Dict]) -> List[List[Dict]]:
        """Group text blocks into rows based on vertical alignment"""
        if not blocks:
            return []
        
        # Sort blocks by vertical position first
        sorted_blocks = sorted(blocks, key=lambda x: x['top'])
        
        rows = []
        current_row = [sorted_blocks[0]]
        current_y = sorted_blocks[0]['top']
        y_threshold = 12  # Tighter threshold for better row detection
        
        for block in sorted_blocks[1:]:
            # Check if this block belongs to the current row
            if abs(block['top'] - current_y) <= y_threshold:
                current_row.append(block)
            else:
                # Start new row
                if current_row:
                    # Sort current row by horizontal position
                    current_row.sort(key=lambda x: x['left'])
                    rows.append(current_row)
                
                current_row = [block]
                current_y = block['top']
        
        # Add the last row
        if current_row:
            current_row.sort(key=lambda x: x['left'])
            rows.append(current_row)
        
        return rows

    def _find_best_column(self, left_pos: int, column_positions: List[int]) -> int:
        """Find the best column index for a given horizontal position"""
        if not column_positions:
            return 0
        
        # Find the closest column position
        best_col = 0
        min_distance = abs(left_pos - column_positions[0])
        
        for i, col_pos in enumerate(column_positions):
            distance = abs(left_pos - col_pos)
            if distance < min_distance:
                min_distance = distance
                best_col = i
        
        return best_col


def main():
    parser = argparse.ArgumentParser(description='Convert PDF to Office documents using OCR')
    parser.add_argument('input_pdf', help='Input PDF file path')
    parser.add_argument('output_file', help='Output file path')
    parser.add_argument('--format', choices=['docx', 'xlsx', 'pptx'], required=True, 
                       help='Output format')
    parser.add_argument('--ocr-engine', choices=['tesseract', 'easyocr'], default='tesseract',
                       help='OCR engine to use')
    parser.add_argument('--dpi', type=int, default=200, help='DPI for PDF to image conversion')
    
    args = parser.parse_args()
    
    try:
        # Initialize converter
        converter = PDFOCRConverter(ocr_engine=args.ocr_engine)
        
        # Extract text from PDF
        pages_data = converter.extract_text_from_pdf(args.input_pdf, dpi=args.dpi)
        
        if not pages_data:
            print("ERROR: No text could be extracted from the PDF")
            sys.exit(1)
        
        # Create output document based on format
        if args.format == 'docx':
            converter.create_word_document(pages_data, args.output_file)
        elif args.format == 'xlsx':
            converter.create_excel_document(pages_data, args.output_file)
        elif args.format == 'pptx':
            converter.create_powerpoint_document(pages_data, args.output_file)
        
        # Return success info as JSON
        total_text = sum(len(page['full_text']) for page in pages_data)
        result = {
            'success': True,
            'pages_processed': len(pages_data),
            'total_characters': total_text,
            'output_file': args.output_file,
            'format': args.format
        }
        
        print(f"\nSUCCESS: {json.dumps(result)}")
        
    except Exception as e:
        error_result = {
            'success': False,
            'error': str(e),
            'traceback': traceback.format_exc()
        }
        print(f"\nERROR: {json.dumps(error_result)}")
        sys.exit(1)


if __name__ == '__main__':
    main()
